#!/usr/bin/env python3
"""Deterministic PPTX creation, inspection, and optional PDF rendering for ArcForge."""

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Mapping, Optional, Sequence, Tuple

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    PPTX_IMPORT_ERROR: Optional[Exception] = None
except (
    Exception
) as error:  # pragma: no cover - exercised only when dependency is absent
    PPTX_IMPORT_ERROR = error


DEFAULT_THEME = {
    "background": "F8FAFC",
    "surface": "E2E8F0",
    "primary": "0F172A",
    "accent": "2563EB",
    "text": "0F172A",
    "muted": "64748B",
    "inverse": "F8FAFC",
    "font": "Microsoft YaHei",
    "font_alt": "Aptos",
}
SERIES_COLORS = ["2563EB", "0F172A", "0891B2", "7C3AED", "EA580C", "16A34A"]
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5


class PresentationError(RuntimeError):
    pass


def require_pptx() -> None:
    if PPTX_IMPORT_ERROR is not None:
        raise PresentationError(
            "python-pptx and Pillow are required. Ask the user before installing "
            "scripts/requirements.txt. Import error: " + str(PPTX_IMPORT_ERROR)
        )


def load_json_object(path_value: str) -> Tuple[Dict[str, Any], Path]:
    path = Path(path_value).expanduser().resolve()
    if not path.is_file():
        raise PresentationError("JSON specification does not exist: " + str(path))
    try:
        with path.open("r", encoding="utf-8-sig") as handle:
            value = json.load(handle)
    except (OSError, json.JSONDecodeError) as error:
        raise PresentationError(
            "Failed to read JSON specification: " + str(error)
        ) from error
    if not isinstance(value, dict):
        raise PresentationError("JSON specification must be an object")
    return value, path


def normalize_color(value: Any, field: str = "color") -> str:
    raw = str(value).strip().lstrip("#").upper()
    if len(raw) != 6 or any(char not in "0123456789ABCDEF" for char in raw):
        raise PresentationError(field + " must be a 6-digit hexadecimal color")
    return raw


def merged_theme(spec: Mapping[str, Any]) -> Dict[str, str]:
    theme = dict(DEFAULT_THEME)
    raw = spec.get("theme")
    if raw is None:
        return theme
    if not isinstance(raw, dict):
        raise PresentationError("theme must be an object")
    for key in (
        "background",
        "surface",
        "primary",
        "accent",
        "text",
        "muted",
        "inverse",
    ):
        if key in raw:
            theme[key] = normalize_color(raw[key], "theme." + key)
    for key in ("font", "font_alt"):
        if key in raw and str(raw[key]).strip():
            theme[key] = str(raw[key]).strip()
    return theme


def rgb(value: str) -> Any:
    return RGBColor.from_string(normalize_color(value))


def set_slide_background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_solid_shape(
    slide: Any,
    shape_type: Any,
    x: float,
    y: float,
    width: float,
    height: float,
    fill_color: str,
    line_color: Optional[str] = None,
    radius_name: Optional[str] = None,
) -> Any:
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
    else:
        shape.line.fill.background()
    if radius_name:
        shape.name = radius_name
    return shape


def add_text(
    slide: Any,
    text: Any,
    x: float,
    y: float,
    width: float,
    height: float,
    theme: Mapping[str, str],
    font_size: float = 18,
    color: Optional[str] = None,
    bold: bool = False,
    align: Any = None,
    vertical: Any = None,
    font_name: Optional[str] = None,
    margin: float = 0.02,
    name: Optional[str] = None,
) -> Any:
    shape = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    if vertical is not None:
        frame.vertical_anchor = vertical
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align if align is not None else PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font_name or theme["font"]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color or theme["text"])
    return shape


def add_content_header(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> float:
    set_slide_background(slide, str(slide_spec.get("background", theme["background"])))
    accent = normalize_color(slide_spec.get("accent", theme["accent"]), "slide.accent")
    add_solid_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.09, accent)
    title = str(slide_spec.get("title", "")).strip()
    if not title:
        raise PresentationError("Content slide requires a title")
    add_text(
        slide,
        title,
        0.7,
        0.38,
        11.9,
        0.62,
        theme,
        font_size=26,
        bold=True,
        name="ArcForge Title",
    )
    subtitle = str(slide_spec.get("subtitle", "")).strip()
    if subtitle:
        add_text(
            slide,
            subtitle,
            0.72,
            1.0,
            11.6,
            0.35,
            theme,
            font_size=11.5,
            color=theme["muted"],
        )
        return 1.52
    return 1.30


def add_footer(
    slide: Any,
    footer_text: str,
    slide_number: int,
    theme: Mapping[str, str],
) -> None:
    add_solid_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        0.7,
        6.94,
        11.93,
        0.012,
        theme["surface"],
    )
    if footer_text:
        add_text(
            slide,
            footer_text,
            0.7,
            7.01,
            10.9,
            0.22,
            theme,
            font_size=8.5,
            color=theme["muted"],
        )
    add_text(
        slide,
        str(slide_number),
        11.75,
        7.01,
        0.85,
        0.22,
        theme,
        font_size=8.5,
        color=theme["muted"],
        align=PP_ALIGN.RIGHT,
    )


def normalized_bullets(raw: Any) -> List[Dict[str, Any]]:
    if not isinstance(raw, list):
        raise PresentationError("bullets must be an array")
    bullets: List[Dict[str, Any]] = []
    for item in raw:
        if isinstance(item, str):
            bullets.append({"text": item, "level": 0})
        elif isinstance(item, dict) and str(item.get("text", "")).strip():
            bullets.append(
                {
                    "text": str(item["text"]),
                    "level": max(0, min(3, int(item.get("level", 0)))),
                    "accent": item.get("accent"),
                }
            )
        else:
            raise PresentationError("each bullet must be text or an object with text")
    return bullets


def add_bullet_list(
    slide: Any,
    bullets: Sequence[Mapping[str, Any]],
    x: float,
    y: float,
    width: float,
    height: float,
    theme: Mapping[str, str],
    font_size: float = 18,
) -> None:
    if not bullets:
        return
    row_height = min(0.72, height / max(1, len(bullets)))
    for index, item in enumerate(bullets):
        level = int(item.get("level", 0))
        item_y = y + index * row_height
        indent = level * 0.30
        marker_color = (
            normalize_color(item["accent"], "bullet.accent")
            if item.get("accent")
            else theme["accent"]
        )
        add_solid_shape(
            slide,
            MSO_SHAPE.OVAL,
            x + indent,
            item_y + 0.18,
            0.10,
            0.10,
            marker_color,
        )
        add_text(
            slide,
            item["text"],
            x + 0.24 + indent,
            item_y,
            width - 0.24 - indent,
            row_height,
            theme,
            font_size=max(12, font_size - level * 1.5),
            vertical=MSO_ANCHOR.MIDDLE,
        )


def render_title_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> None:
    set_slide_background(slide, str(slide_spec.get("background", theme["primary"])))
    accent = normalize_color(slide_spec.get("accent", theme["accent"]), "slide.accent")
    add_solid_shape(slide, MSO_SHAPE.RECTANGLE, 0.78, 1.22, 0.12, 4.75, accent)
    kicker = str(slide_spec.get("kicker", "")).strip()
    if kicker:
        add_text(
            slide,
            kicker.upper(),
            1.18,
            1.18,
            10.9,
            0.38,
            theme,
            font_size=10,
            color=accent,
            bold=True,
        )
    title = str(slide_spec.get("title", "")).strip()
    if not title:
        raise PresentationError("title slide requires a title")
    add_text(
        slide,
        title,
        1.16,
        1.72,
        10.9,
        2.0,
        theme,
        font_size=34,
        color=theme["inverse"],
        bold=True,
        vertical=MSO_ANCHOR.MIDDLE,
        name="ArcForge Title",
    )
    subtitle = str(slide_spec.get("subtitle", "")).strip()
    if subtitle:
        add_text(
            slide,
            subtitle,
            1.18,
            4.10,
            9.9,
            1.0,
            theme,
            font_size=17,
            color="CBD5E1",
        )
    footer = str(slide_spec.get("footer", "ArcForge")).strip()
    add_text(
        slide,
        footer,
        1.18,
        6.67,
        10.0,
        0.32,
        theme,
        font_size=9,
        color="94A3B8",
    )


def render_section_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> None:
    set_slide_background(slide, str(slide_spec.get("background", theme["primary"])))
    accent = normalize_color(slide_spec.get("accent", theme["accent"]), "slide.accent")
    number = str(slide_spec.get("number", "")).strip()
    if number:
        add_text(
            slide,
            number,
            0.85,
            0.62,
            2.3,
            1.3,
            theme,
            font_size=58,
            color=accent,
            bold=True,
        )
    title = str(slide_spec.get("title", "")).strip()
    if not title:
        raise PresentationError("section slide requires a title")
    add_text(
        slide,
        title,
        0.88,
        2.18,
        11.2,
        1.45,
        theme,
        font_size=34,
        color=theme["inverse"],
        bold=True,
        vertical=MSO_ANCHOR.MIDDLE,
        name="ArcForge Title",
    )
    subtitle = str(slide_spec.get("subtitle", "")).strip()
    if subtitle:
        add_text(
            slide,
            subtitle,
            0.90,
            3.92,
            10.8,
            0.75,
            theme,
            font_size=16,
            color="CBD5E1",
        )
    add_solid_shape(slide, MSO_SHAPE.RECTANGLE, 0.88, 5.22, 2.2, 0.08, accent)


def render_bullets_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> None:
    content_top = add_content_header(slide, slide_spec, theme)
    bullets = normalized_bullets(slide_spec.get("bullets", []))
    if len(bullets) > 8:
        raise PresentationError(
            "bullets slide supports at most 8 bullets; split dense content"
        )
    add_bullet_list(slide, bullets, 0.92, content_top + 0.12, 11.25, 5.0, theme, 18)


def render_two_column_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> None:
    content_top = add_content_header(slide, slide_spec, theme)
    for index, key in enumerate(("left", "right")):
        column = slide_spec.get(key, {})
        if not isinstance(column, dict):
            raise PresentationError(key + " must be an object")
        x = 0.70 if index == 0 else 6.72
        add_solid_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            content_top + 0.08,
            5.90,
            4.92,
            "FFFFFF",
            line_color=theme["surface"],
        )
        heading = str(column.get("heading", key.title())).strip()
        add_text(
            slide,
            heading,
            x + 0.38,
            content_top + 0.35,
            5.10,
            0.48,
            theme,
            font_size=17,
            bold=True,
        )
        bullets = normalized_bullets(column.get("bullets", []))
        if len(bullets) > 6:
            raise PresentationError(key + " column supports at most 6 bullets")
        add_bullet_list(
            slide,
            bullets,
            x + 0.40,
            content_top + 1.02,
            5.00,
            3.55,
            theme,
            14,
        )


def render_metrics_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> None:
    content_top = add_content_header(slide, slide_spec, theme)
    metrics = slide_spec.get("metrics")
    if not isinstance(metrics, list) or not metrics:
        raise PresentationError("metrics slide requires a non-empty metrics array")
    if len(metrics) > 4:
        raise PresentationError("metrics slide supports at most 4 metric cards")
    gap = 0.24
    total_width = 11.93
    card_width = (total_width - gap * (len(metrics) - 1)) / len(metrics)
    for index, metric in enumerate(metrics):
        if not isinstance(metric, dict):
            raise PresentationError("each metric must be an object")
        x = 0.70 + index * (card_width + gap)
        add_solid_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            content_top + 0.35,
            card_width,
            3.72,
            "FFFFFF",
            line_color=theme["surface"],
        )
        add_solid_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            x,
            content_top + 0.35,
            card_width,
            0.09,
            theme["accent"],
        )
        add_text(
            slide,
            metric.get("value", ""),
            x + 0.28,
            content_top + 0.90,
            card_width - 0.56,
            0.90,
            theme,
            font_size=28,
            bold=True,
            vertical=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            metric.get("label", ""),
            x + 0.28,
            content_top + 1.92,
            card_width - 0.56,
            0.90,
            theme,
            font_size=12,
            color=theme["muted"],
        )
        if metric.get("delta") is not None:
            add_text(
                slide,
                metric["delta"],
                x + 0.28,
                content_top + 3.08,
                card_width - 0.56,
                0.34,
                theme,
                font_size=11,
                color=theme["accent"],
                bold=True,
            )


def set_table_cell(
    cell: Any,
    value: Any,
    theme: Mapping[str, str],
    fill_color: str,
    text_color: str,
    bold: bool,
    font_size: float,
) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(fill_color)
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(value)
    run.font.name = theme["font"]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(text_color)


def render_table_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> None:
    content_top = add_content_header(slide, slide_spec, theme)
    columns = slide_spec.get("columns")
    rows = slide_spec.get("rows")
    if not isinstance(columns, list) or not columns:
        raise PresentationError("table slide requires a non-empty columns array")
    if not isinstance(rows, list):
        raise PresentationError("table slide rows must be an array")
    if len(columns) > 8 or len(rows) > 12:
        raise PresentationError(
            "table exceeds the supported 8 columns or 12 rows; split it"
        )
    total_rows = len(rows) + 1
    x, y, width, height = 0.70, content_top + 0.18, 11.93, 5.12
    table_shape = slide.shapes.add_table(
        total_rows,
        len(columns),
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    table_shape.name = "ArcForge Table"
    table = table_shape.table

    raw_widths = slide_spec.get("column_widths")
    if raw_widths is not None:
        if not isinstance(raw_widths, list) or len(raw_widths) != len(columns):
            raise PresentationError("column_widths must match the number of columns")
        weights = [float(item) for item in raw_widths]
        if any(item <= 0 for item in weights):
            raise PresentationError("column_widths values must be positive")
    else:
        weights = [1.0] * len(columns)
    weight_total = sum(weights)
    for index, weight in enumerate(weights):
        table.columns[index].width = Inches(width * weight / weight_total)

    for column_index, value in enumerate(columns):
        set_table_cell(
            table.cell(0, column_index),
            value,
            theme,
            theme["primary"],
            theme["inverse"],
            True,
            11,
        )
    body_font_size = 10.5 if len(rows) <= 8 else 9
    for row_index, raw_row in enumerate(rows, start=1):
        if not isinstance(raw_row, list):
            raise PresentationError("each table row must be an array")
        values = list(raw_row[: len(columns)]) + [""] * max(
            0, len(columns) - len(raw_row)
        )
        fill_color = "FFFFFF" if row_index % 2 else "F1F5F9"
        for column_index, value in enumerate(values):
            set_table_cell(
                table.cell(row_index, column_index),
                value,
                theme,
                fill_color,
                theme["text"],
                False,
                body_font_size,
            )


def chart_type_value(value: str) -> Any:
    normalized = value.lower()
    mapping = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    if normalized not in mapping:
        raise PresentationError("Unsupported chart_type: " + value)
    return mapping[normalized]


def render_chart_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> None:
    content_top = add_content_header(slide, slide_spec, theme)
    categories = slide_spec.get("categories")
    series = slide_spec.get("series")
    if not isinstance(categories, list) or not categories:
        raise PresentationError("chart slide requires a non-empty categories array")
    if not isinstance(series, list) or not series:
        raise PresentationError("chart slide requires a non-empty series array")

    chart_data = ChartData()
    chart_data.categories = [str(item) for item in categories]
    for item in series:
        if not isinstance(item, dict) or not str(item.get("name", "")).strip():
            raise PresentationError("each chart series requires a name and values")
        values = item.get("values")
        if not isinstance(values, list) or len(values) != len(categories):
            raise PresentationError(
                "each chart series values array must match categories"
            )
        chart_data.add_series(str(item["name"]), [float(value) for value in values])

    chart_type = str(slide_spec.get("chart_type", "column"))
    chart_frame = slide.shapes.add_chart(
        chart_type_value(chart_type),
        Inches(0.82),
        Inches(content_top + 0.12),
        Inches(11.68),
        Inches(5.08),
        chart_data,
    )
    chart_frame.name = "ArcForge Chart"
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.chart_style = int(slide_spec.get("chart_style", 10))

    for index, chart_series in enumerate(chart.series):
        color = rgb(SERIES_COLORS[index % len(SERIES_COLORS)])
        if chart_type.lower() == "line":
            chart_series.format.line.color.rgb = color
            chart_series.format.line.width = Pt(2.25)
        else:
            chart_series.format.fill.solid()
            chart_series.format.fill.fore_color.rgb = color
            chart_series.format.line.color.rgb = color

    if chart_type.lower() not in ("pie", "doughnut"):
        chart.category_axis.tick_labels.font.name = theme["font_alt"]
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.name = theme["font_alt"]
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = rgb(theme["surface"])
    else:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_percentage = True
        plot.data_labels.show_legend_key = False


def resolve_asset_path(raw_path: Any, spec_dir: Path) -> Path:
    path = Path(str(raw_path)).expanduser()
    if not path.is_absolute():
        path = spec_dir / path
    path = path.resolve()
    if not path.is_file():
        raise PresentationError("Image does not exist: " + str(path))
    return path


def add_picture_fit(
    slide: Any,
    image_path: Path,
    x: float,
    y: float,
    width: float,
    height: float,
    fit: str,
) -> Any:
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    if image_width <= 0 or image_height <= 0:
        raise PresentationError("Image has invalid dimensions: " + str(image_path))
    image_ratio = image_width / image_height
    frame_ratio = width / height

    if fit == "contain":
        if image_ratio >= frame_ratio:
            fitted_width = width
            fitted_height = width / image_ratio
        else:
            fitted_height = height
            fitted_width = height * image_ratio
        left = x + (width - fitted_width) / 2
        top = y + (height - fitted_height) / 2
        picture = slide.shapes.add_picture(
            str(image_path),
            Inches(left),
            Inches(top),
            Inches(fitted_width),
            Inches(fitted_height),
        )
    elif fit == "cover":
        picture = slide.shapes.add_picture(
            str(image_path),
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(height),
        )
        if image_ratio > frame_ratio:
            crop = (1.0 - frame_ratio / image_ratio) / 2.0
            picture.crop_left = crop
            picture.crop_right = crop
        elif image_ratio < frame_ratio:
            crop = (1.0 - image_ratio / frame_ratio) / 2.0
            picture.crop_top = crop
            picture.crop_bottom = crop
    else:
        raise PresentationError("image fit must be contain or cover")
    picture.name = "ArcForge Image"
    return picture


def render_image_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
    spec_dir: Path,
) -> None:
    content_top = add_content_header(slide, slide_spec, theme)
    if not slide_spec.get("image"):
        raise PresentationError("image slide requires an image path")
    image_path = resolve_asset_path(slide_spec["image"], spec_dir)
    caption = str(slide_spec.get("caption", "")).strip()
    box_height = 4.55 if caption else 5.12
    add_solid_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.70,
        content_top + 0.10,
        11.93,
        box_height,
        "FFFFFF",
        line_color=theme["surface"],
    )
    add_picture_fit(
        slide,
        image_path,
        0.84,
        content_top + 0.24,
        11.65,
        box_height - 0.28,
        str(slide_spec.get("fit", "contain")).lower(),
    )
    if caption:
        add_text(
            slide,
            caption,
            0.82,
            content_top + 4.83,
            11.55,
            0.34,
            theme,
            font_size=9.5,
            color=theme["muted"],
            align=PP_ALIGN.CENTER,
        )


def render_quote_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> None:
    set_slide_background(slide, str(slide_spec.get("background", theme["background"])))
    quote = str(slide_spec.get("quote", "")).strip()
    if not quote:
        raise PresentationError("quote slide requires quote text")
    title = str(slide_spec.get("title", "")).strip()
    if title:
        add_text(
            slide,
            title,
            0.72,
            0.45,
            11.7,
            0.55,
            theme,
            font_size=18,
            color=theme["muted"],
            name="ArcForge Title",
        )
    add_text(
        slide,
        "“",
        0.82,
        1.30,
        1.0,
        1.0,
        theme,
        font_size=68,
        color=theme["accent"],
        bold=True,
    )
    add_text(
        slide,
        quote,
        1.48,
        1.72,
        10.35,
        3.10,
        theme,
        font_size=27,
        bold=True,
        vertical=MSO_ANCHOR.MIDDLE,
        name=None if title else "ArcForge Title",
    )
    attribution = str(slide_spec.get("attribution", "")).strip()
    if attribution:
        add_text(
            slide,
            "— " + attribution,
            1.52,
            5.12,
            9.6,
            0.55,
            theme,
            font_size=13,
            color=theme["muted"],
        )


def render_closing_slide(
    slide: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
) -> None:
    set_slide_background(slide, str(slide_spec.get("background", theme["primary"])))
    accent = normalize_color(slide_spec.get("accent", theme["accent"]), "slide.accent")
    add_solid_shape(slide, MSO_SHAPE.RECTANGLE, 0.82, 1.25, 0.10, 4.60, accent)
    title = str(slide_spec.get("title", "")).strip()
    if not title:
        raise PresentationError("closing slide requires a title")
    add_text(
        slide,
        title,
        1.22,
        1.65,
        10.6,
        1.65,
        theme,
        font_size=31,
        color=theme["inverse"],
        bold=True,
        vertical=MSO_ANCHOR.MIDDLE,
        name="ArcForge Title",
    )
    subtitle = str(slide_spec.get("subtitle", "")).strip()
    if subtitle:
        add_text(
            slide,
            subtitle,
            1.24,
            3.72,
            9.9,
            0.92,
            theme,
            font_size=16,
            color="CBD5E1",
        )
    contact = str(slide_spec.get("contact", "")).strip()
    if contact:
        add_text(
            slide,
            contact,
            1.24,
            5.42,
            9.8,
            0.42,
            theme,
            font_size=11,
            color=accent,
            bold=True,
        )


def apply_notes(slide: Any, notes: Any) -> None:
    if notes is None or not str(notes).strip():
        return
    try:
        slide.notes_slide.notes_text_frame.text = str(notes)
    except Exception as error:
        raise PresentationError("Failed to write slide notes: " + str(error)) from error


def render_slide(
    presentation: Any,
    slide_spec: Mapping[str, Any],
    theme: Mapping[str, str],
    spec_dir: Path,
    footer_text: str,
    slide_number: int,
) -> None:
    if not isinstance(slide_spec, dict):
        raise PresentationError("each slides entry must be an object")
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide_type = str(slide_spec.get("type", "bullets")).strip().lower()
    if slide_type == "title":
        render_title_slide(slide, slide_spec, theme)
    elif slide_type == "section":
        render_section_slide(slide, slide_spec, theme)
    elif slide_type == "bullets":
        render_bullets_slide(slide, slide_spec, theme)
    elif slide_type == "two-column":
        render_two_column_slide(slide, slide_spec, theme)
    elif slide_type == "metrics":
        render_metrics_slide(slide, slide_spec, theme)
    elif slide_type == "table":
        render_table_slide(slide, slide_spec, theme)
    elif slide_type == "chart":
        render_chart_slide(slide, slide_spec, theme)
    elif slide_type == "image":
        render_image_slide(slide, slide_spec, theme, spec_dir)
    elif slide_type == "quote":
        render_quote_slide(slide, slide_spec, theme)
    elif slide_type == "closing":
        render_closing_slide(slide, slide_spec, theme)
    else:
        raise PresentationError("Unsupported slide type: " + slide_type)
    if slide_type not in ("title", "section", "closing"):
        add_footer(slide, footer_text, slide_number, theme)
    apply_notes(slide, slide_spec.get("notes"))


def apply_metadata(presentation: Any, spec: Mapping[str, Any]) -> None:
    metadata = spec.get("metadata", {})
    if not isinstance(metadata, dict):
        raise PresentationError("metadata must be an object")
    properties = presentation.core_properties
    mapping = {
        "title": "title",
        "subject": "subject",
        "author": "author",
        "comments": "comments",
        "keywords": "keywords",
        "category": "category",
    }
    for source, target in mapping.items():
        if source in metadata:
            setattr(properties, target, str(metadata[source]))
    if not properties.author:
        properties.author = "ArcForge"
    if metadata.get("company") and not properties.comments:
        properties.comments = "Company: " + str(metadata["company"])


def create_presentation(spec: Mapping[str, Any], spec_dir: Path) -> Any:
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        raise PresentationError("slides must be a non-empty array")
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH)
    presentation.slide_height = Inches(SLIDE_HEIGHT)
    apply_metadata(presentation, spec)
    theme = merged_theme(spec)
    footer = str(spec.get("footer", "")).strip()
    for index, slide_spec in enumerate(slides, start=1):
        render_slide(presentation, slide_spec, theme, spec_dir, footer, index)
    return presentation


def normalized_output_path(path_value: str, suffix: str) -> Path:
    path = Path(path_value).expanduser().resolve()
    if path.suffix.lower() != suffix:
        raise PresentationError("Output path must end with " + suffix)
    return path


def atomic_save(presentation: Any, output_path: Path, force: bool) -> None:
    if output_path.exists() and not force:
        raise PresentationError(
            "Output already exists. Use a new path or pass --force only after explicit approval: "
            + str(output_path)
        )
    output_path.parent.mkdir(parents=True, exist_ok=True)
    handle, temporary_name = tempfile.mkstemp(
        prefix=".arcforge-deck-",
        suffix=".pptx",
        dir=str(output_path.parent),
    )
    os.close(handle)
    temporary_path = Path(temporary_name)
    try:
        presentation.save(temporary_path)
        Presentation(temporary_path)
        os.replace(temporary_path, output_path)
    except Exception:
        try:
            temporary_path.unlink(missing_ok=True)
        except OSError:
            pass
        raise


def shape_is_out_of_bounds(shape: Any, slide_width: int, slide_height: int) -> bool:
    return (
        shape.left < 0
        or shape.top < 0
        or shape.left + shape.width > slide_width
        or shape.top + shape.height > slide_height
    )


def slide_title(slide: Any) -> str:
    for shape in slide.shapes:
        if shape.name == "ArcForge Title" and getattr(shape, "has_text_frame", False):
            return shape.text.strip()
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            return shape.text.strip().splitlines()[0]
    return ""


def inspect_presentation(path: Path) -> Dict[str, Any]:
    if not path.is_file():
        raise PresentationError("Presentation does not exist: " + str(path))
    try:
        presentation = Presentation(path)
    except Exception as error:
        raise PresentationError(
            "Failed to inspect presentation: " + str(error)
        ) from error

    slides: List[Dict[str, Any]] = []
    total_images = 0
    total_tables = 0
    total_charts = 0
    total_out_of_bounds = 0
    for index, slide in enumerate(presentation.slides, start=1):
        images = 0
        tables = 0
        charts = 0
        out_of_bounds = 0
        text_characters = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images += 1
            if getattr(shape, "has_table", False):
                tables += 1
            if getattr(shape, "has_chart", False):
                charts += 1
            if getattr(shape, "has_text_frame", False):
                text_characters += len(shape.text)
            if shape_is_out_of_bounds(
                shape, presentation.slide_width, presentation.slide_height
            ):
                out_of_bounds += 1
        title = slide_title(slide)
        total_images += images
        total_tables += tables
        total_charts += charts
        total_out_of_bounds += out_of_bounds
        notes = ""
        if slide.has_notes_slide:
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes = ""
        slides.append(
            {
                "number": index,
                "title": title,
                "shape_count": len(slide.shapes),
                "text_characters": text_characters,
                "image_count": images,
                "table_count": tables,
                "chart_count": charts,
                "out_of_bounds_shapes": out_of_bounds,
                "has_notes": bool(notes),
            }
        )

    return {
        "path": str(path.resolve()),
        "size_bytes": path.stat().st_size,
        "slide_count": len(slides),
        "slide_size_inches": {
            "width": round(presentation.slide_width / Inches(1), 3),
            "height": round(presentation.slide_height / Inches(1), 3),
        },
        "slides": slides,
        "total_images": total_images,
        "total_tables": total_tables,
        "total_charts": total_charts,
        "total_out_of_bounds_shapes": total_out_of_bounds,
        "missing_title_slides": [
            item["number"] for item in slides if not item["title"]
        ],
        "visually_rendered": False,
    }


def soffice_candidates() -> List[Path]:
    candidates: List[Path] = []
    for key in ("ARCFORGE_SOFFICE_PATH", "LIVEAGENT_SOFFICE_PATH"):
        raw = os.environ.get(key)
        if raw:
            candidates.append(Path(raw).expanduser())
    for command in ("soffice", "libreoffice"):
        resolved = shutil.which(command)
        if resolved:
            candidates.append(Path(resolved))
    for key in ("ProgramFiles", "ProgramFiles(x86)", "LOCALAPPDATA"):
        root = os.environ.get(key)
        if root:
            candidates.append(Path(root) / "LibreOffice" / "program" / "soffice.exe")
    candidates.extend(
        [
            Path("/usr/bin/soffice"),
            Path("/usr/local/bin/soffice"),
            Path("/opt/libreoffice/program/soffice"),
            Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"),
        ]
    )
    return candidates


def find_soffice() -> Path:
    for candidate in soffice_candidates():
        try:
            if candidate.is_file():
                return candidate.resolve()
        except OSError:
            continue
    raise PresentationError(
        "LibreOffice soffice was not found. Set ARCFORGE_SOFFICE_PATH or install "
        "LibreOffice only after user approval."
    )


def render_pdf(input_path: Path, output_path: Path, force: bool) -> Dict[str, Any]:
    if not input_path.is_file():
        raise PresentationError("Input presentation does not exist: " + str(input_path))
    if output_path.exists() and not force:
        raise PresentationError(
            "Output already exists. Use a new path or pass --force only after explicit approval: "
            + str(output_path)
        )
    output_path.parent.mkdir(parents=True, exist_ok=True)
    soffice = find_soffice()
    with tempfile.TemporaryDirectory(prefix="arcforge-slides-") as temporary_dir:
        temporary_root = Path(temporary_dir)
        profile = temporary_root / "profile"
        profile.mkdir()
        command = [
            str(soffice),
            "--headless",
            "-env:UserInstallation=" + profile.resolve().as_uri(),
            "--convert-to",
            "pdf",
            "--outdir",
            str(temporary_root),
            str(input_path),
        ]
        try:
            completed = subprocess.run(
                command,
                check=False,
                capture_output=True,
                text=True,
                timeout=180,
            )
        except (OSError, subprocess.TimeoutExpired) as error:
            raise PresentationError(
                "Failed to run LibreOffice: " + str(error)
            ) from error
        if completed.returncode != 0:
            detail = (completed.stderr or completed.stdout).strip()
            raise PresentationError("LibreOffice PDF conversion failed: " + detail)
        generated = temporary_root / (input_path.stem + ".pdf")
        if not generated.is_file():
            raise PresentationError("LibreOffice did not produce the expected PDF")
        os.replace(generated, output_path)
    return {
        "path": str(output_path),
        "size_bytes": output_path.stat().st_size,
        "renderer": str(soffice),
        "visually_rendered": True,
    }


def run_create(args: argparse.Namespace) -> Dict[str, Any]:
    spec, spec_path = load_json_object(args.spec)
    output_path = normalized_output_path(args.output, ".pptx")
    presentation = create_presentation(spec, spec_path.parent)
    atomic_save(presentation, output_path, args.force)
    return {"action": "created", "presentation": inspect_presentation(output_path)}


def run_inspect(args: argparse.Namespace) -> Dict[str, Any]:
    input_path = Path(args.input).expanduser().resolve()
    return {"action": "inspected", "presentation": inspect_presentation(input_path)}


def run_render(args: argparse.Namespace) -> Dict[str, Any]:
    input_path = Path(args.input).expanduser().resolve()
    output_path = normalized_output_path(args.output, ".pdf")
    return {
        "action": "rendered",
        "input": str(input_path),
        "pdf": render_pdf(input_path, output_path, args.force),
    }


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Create, inspect, and optionally render PowerPoint decks for ArcForge."
    )
    subparsers = parser.add_subparsers(dest="command", required=True)

    create_parser = subparsers.add_parser("create", help="Create a PPTX deck from JSON")
    create_parser.add_argument("--spec", required=True, help="UTF-8 JSON specification")
    create_parser.add_argument("--output", required=True, help="Destination .pptx path")
    create_parser.add_argument(
        "--force", action="store_true", help="Overwrite the exact output path"
    )
    create_parser.set_defaults(handler=run_create)

    inspect_parser = subparsers.add_parser(
        "inspect", help="Print a structural deck summary"
    )
    inspect_parser.add_argument("--input", required=True, help="Existing .pptx path")
    inspect_parser.set_defaults(handler=run_inspect)

    render_parser = subparsers.add_parser(
        "render", help="Render a PPTX deck to PDF with LibreOffice"
    )
    render_parser.add_argument("--input", required=True, help="Existing .pptx path")
    render_parser.add_argument("--output", required=True, help="Destination .pdf path")
    render_parser.add_argument(
        "--force", action="store_true", help="Overwrite the exact output path"
    )
    render_parser.set_defaults(handler=run_render)
    return parser


def main(argv: Optional[Sequence[str]] = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        require_pptx()
        result = args.handler(args)
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0
    except PresentationError as error:
        print("error: " + str(error), file=sys.stderr)
        return 2
    except Exception as error:
        print("error: unexpected presentation failure: " + str(error), file=sys.stderr)
        return 3


if __name__ == "__main__":
    raise SystemExit(main())
